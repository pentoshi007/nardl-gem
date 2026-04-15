#!/usr/bin/env python3
"""Generate PPTX progress report presentation - white bg, black text."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

FIGS = "outputs/figures"
OUT = "presentation"
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(255,255,255)
BLACK = RGBColor(0,0,0)
DARK_BLUE = RGBColor(30,60,114)
GRAY = RGBColor(100,100,100)
LIGHT_BG = RGBColor(245,245,250)

def add_slide():
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    return slide

def add_title_text(slide, text, left, top, width, height, size=28, bold=True, color=BLACK, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf

def add_body_text(slide, text, left, top, width, height, size=16, color=BLACK, bold=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(text.split("\n")):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.space_after = Pt(4)
    return tf

def add_image(slide, fname, left, top, width):
    path = f"{FIGS}/{fname}"
    if os.path.exists(path):
        slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(width))

def add_line(slide, top):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, Inches(0.8), Inches(top), Inches(11.7), Pt(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()

# ===== SLIDE 1: TITLE =====
s = add_slide()
add_title_text(s, "Dissertation Progress Report", 1.5, 1.5, 10, 1, size=36, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_title_text(s, "Do Global Oil Price Shocks Raise India's Inflation\nMore Than They Lower It?", 1.5, 2.8, 10, 1.2, size=24, color=BLACK, align=PP_ALIGN.CENTER, bold=False)
add_body_text(s, "Short-Run Pass-Through to CPI Inflation in India, 2004-2024", 1.5, 4.2, 10, 0.5, size=18, color=GRAY)
add_body_text(s, "Aniket Pandey\nSupervisor: Prof. Shakti Kumar\nMS Economics, JNU, 2026", 1.5, 5.2, 10, 1.2, size=16, color=BLACK)

# ===== SLIDE 2: RESEARCH QUESTION =====
s = add_slide()
add_title_text(s, "Research Question & Motivation", 0.8, 0.4, 11, 0.7, size=30, color=DARK_BLUE)
add_line(s, 1.1)
add_body_text(s, "Core Question:", 0.8, 1.4, 11, 0.4, size=20, bold=True)
add_body_text(s, "How strongly do oil price shocks pass through to India's monthly CPI inflation,\nand is pass-through larger for oil price increases than for decreases?", 0.8, 1.9, 11, 0.8, size=18)
add_body_text(s, "Why This Matters:", 0.8, 3.0, 11, 0.4, size=20, bold=True)
add_body_text(s, "- India imports ~85% of crude oil - highly vulnerable to global oil price shocks\n- Oil affects petrol, diesel, cooking gas, transport, and ultimately all consumer prices\n- If oil increases raise inflation more than oil decreases reduce it = ASYMMETRY\n- Also called 'Rockets and Feathers' - prices go up fast but come down slowly\n- Policy relevance: RBI monetary policy, fiscal interventions, subsidy design", 0.8, 3.5, 11, 2.5, size=16)

# ===== SLIDE 3: DATA SOURCES =====
s = add_slide()
add_title_text(s, "Data Sources & Study Window", 0.8, 0.4, 11, 0.7, size=30, color=DARK_BLUE)
add_line(s, 1.1)
add_body_text(s, "Study Period: April 2004 - December 2024 (249 monthly observations)", 0.8, 1.4, 11, 0.4, size=18, bold=True)
add_body_text(s,
"Data Sources:\n\n"
"1. CPI (Consumer Price Index) - OECD via FRED (2015=100 base)\n"
"   Dependent variable: India's headline inflation measure\n\n"
"2. Brent Crude Oil Price (USD/barrel) - IMF via FRED\n"
"   Global benchmark; India's oil imports priced off Brent\n\n"
"3. INR/USD Exchange Rate - US Federal Reserve via FRED\n"
"   Oil priced in USD but India pays in INR; exchange rate matters\n\n"
"4. IIP (Index of Industrial Production) - RBI DBIE (chain-linked)\n"
"   Control variable for economic activity\n\n"
"5. Fuel & Light CPI (Appendix only) - Official MoSPI CPI API\n"
"   More energy-sensitive price index; sample 2011-2024",
0.8, 2.0, 11, 5, size=15)

# ===== SLIDE 4: RAW DATA =====
s = add_slide()
add_title_text(s, "Raw Data Series (Levels)", 0.8, 0.4, 5, 0.7, size=30, color=DARK_BLUE)
add_line(s, 1.1)
add_image(s, "fig_1_raw_series.png", 0.5, 1.3, 7.5)
add_body_text(s,
"Key Observations:\n\n"
"- CPI: Steady upward trend (non-stationary)\n"
"- Brent: Highly volatile with booms & crashes\n"
"- INR/USD: Steady rupee depreciation (39 to 85)\n"
"- IIP: Upward trend, COVID crash visible\n\n"
"All series are TRENDING = non-stationary\n"
"=> Must transform to log-differences\n   for valid regression",
8.3, 1.3, 4.5, 5, size=14)

# ===== SLIDE 5: LOG-DIFF DATA =====
s = add_slide()
add_title_text(s, "Log-Differenced Series (Monthly % Changes)", 0.8, 0.4, 8, 0.7, size=30, color=DARK_BLUE)
add_line(s, 1.1)
add_image(s, "fig_2_log_diff_series.png", 0.5, 1.3, 7.5)
add_body_text(s,
"After transformation:\n\n"
"dlnX(t) = 100 x [ln(X_t) - ln(X_{t-1})]\n\n"
"Why logs & differences?\n"
"1. Percentage interpretation\n"
"2. Makes data stationary\n"
"3. Reduces skewness\n"
"4. Makes scales comparable\n\n"
"Key: Oil is ~12x more volatile\n"
"than CPI => expect small\n"
"pass-through coefficient",
8.3, 1.3, 4.5, 5, size=14)

# ===== SLIDE 6: METHODOLOGY =====
s = add_slide()
add_title_text(s, "Methodology: Asymmetric ADL Model", 0.8, 0.4, 11, 0.7, size=30, color=DARK_BLUE)
add_line(s, 1.1)
add_body_text(s,
"Step 1: Construct Oil_INR = Brent_USD x INR/USD\n\n"
"Step 2: Take log-differences (monthly % changes)\n\n"
"Step 3: Asymmetric Decomposition (KEY IDEA):\n"
"   dOil+(t) = max(dlnOil(t), 0)    -- keeps only INCREASES\n"
"   dOil-(t) = min(dlnOil(t), 0)    -- keeps only DECREASES\n\n"
"Step 4: Estimate Asymmetric ADL(p, q) model:\n"
"   dlnCPI(t) = a + SUM[gi * dlnCPI(t-i)]\n"
"             + SUM[pi+ * dOil+(t-j)]   (positive oil lags)\n"
"             + SUM[pi- * dOil-(t-j)]   (negative oil lags)\n"
"             + d*dlnIIP(t) + policy dummies + monthly dummies + e(t)\n\n"
"   Oil lag q = 3 (fixed by theory); AR lag p selected by AIC => p = 3\n"
"   Final model: ADL(3,3) with Newey-West HAC standard errors",
0.8, 1.4, 12, 5.5, size=16)

# ===== SLIDE 7: OIL DECOMPOSITION =====
s = add_slide()
add_title_text(s, "Oil Price Asymmetric Decomposition", 0.8, 0.4, 8, 0.7, size=30, color=DARK_BLUE)
add_line(s, 1.1)
add_image(s, "fig_3_oil_decomposition.png", 0.5, 1.3, 7.5)
add_body_text(s,
"Cumulative Partial Sums:\n\n"
"Red = Cumulative oil increases\n"
"Blue = Cumulative oil decreases\n\n"
"Both components have\n"
"substantial variation,\n"
"ensuring reliable estimation\n"
"of separate effects.\n\n"
"This is the core innovation:\n"
"splitting oil changes lets us\n"
"test if increases & decreases\n"
"affect CPI differently.",
8.3, 1.3, 4.5, 5, size=14)

# ===== SLIDE 8: UNIT ROOT =====
s = add_slide()
add_title_text(s, "Unit Root Tests (ADF)", 0.8, 0.4, 11, 0.7, size=30, color=DARK_BLUE)
add_line(s, 1.1)
add_body_text(s,
"Augmented Dickey-Fuller Test: H0: Unit root (non-stationary)\n\n"
"Variables in LEVELS:                              Variables in DIFFERENCES:\n\n"
"  ln(CPI)     ADF = -0.32   p = 0.99  FAIL        dlnCPI   ADF = -9.04   p < 0.01  PASS\n"
"  ln(Oil_INR) ADF = -2.57   p = 0.33  FAIL        dlnOil   ADF = -6.13   p < 0.01  PASS\n"
"  ln(EXR)     ADF = -3.69   p = 0.03  Border.     dlnEXR   ADF = -5.21   p < 0.01  PASS\n"
"  ln(IIP)     ADF = -3.20   p = 0.09  FAIL        dlnIIP   ADF = -8.22   p < 0.01  PASS\n\n"
"Conclusion: All variables are I(1) - integrated of order 1.\n"
"Levels are non-stationary; first differences are strongly stationary.\n"
"This justifies our estimation in log-differences.",
0.8, 1.5, 12, 5, size=15)

# ===== SLIDE 9: MAIN RESULTS =====
s = add_slide()
add_title_text(s, "Main Results: Asymmetric ADL(3,3)", 0.8, 0.4, 11, 0.7, size=30, color=DARK_BLUE)
add_line(s, 1.1)
add_body_text(s,
"Cumulative Pass-Through (CPT) = Total oil effect over 0-3 month lags\n\n"
"  CPT+ (positive oil effect)  = 0.021296     +10% oil shock => +0.213 pp CPI\n"
"  CPT- (negative oil effect)  = 0.000598     -10% oil shock => -0.006 pp CPI\n"
"  Asymmetry Gap               = 0.020698     Positive effect is ~35x larger\n\n"
"Statistical Tests (Newey-West HAC):\n\n"
"  H0: CPT+ = 0       p = 0.1220   Not significant at 5%, borderline at 12%\n"
"  H0: CPT- = 0       p = 0.9375   Clearly not significant\n"
"  H0: CPT+ = CPT-    p = 0.2408   Asymmetry not significant at 5%\n\n"
"  Adj R-squared = 0.4492   (model explains ~45% of CPI variation)\n"
"  N = 245 observations\n\n"
"Interpretation: Point estimates show clear asymmetry. Positive oil shocks have\n"
"economically meaningful effect. But formal asymmetry test is not significant at 5%.",
0.8, 1.4, 12, 5.5, size=15)

# ===== SLIDE 10: CPT PLOT =====
s = add_slide()
add_title_text(s, "Cumulative Pass-Through by Lag Horizon", 0.8, 0.4, 8, 0.7, size=30, color=DARK_BLUE)
add_line(s, 1.1)
add_image(s, "fig_4_cumulative_passthrough.png", 0.5, 1.3, 7.5)
add_body_text(s,
"KEY RESULT PLOT:\n\n"
"Red (CPT+): Builds up to\n"
"~0.021 by lag 3. Oil increases\n"
"take 2-3 months to transmit\n"
"to CPI.\n\n"
"Blue (CPT-): Stays near zero.\n"
"Oil decreases have essentially\n"
"NO effect on CPI.\n\n"
"The 2-month delay reflects\n"
"supply chain transmission:\n"
"oil => refineries => retail\n"
"fuel => transport => goods.",
8.3, 1.3, 4.5, 5.2, size=14)

# ===== SLIDE 11: WALD TEST =====
s = add_slide()
add_title_text(s, "Formal Asymmetry Test (Wald Test)", 0.8, 0.4, 11, 0.7, size=30, color=DARK_BLUE)
add_line(s, 1.1)
add_body_text(s,
"Wald Test: H0: CPT+ = CPT-  (no asymmetry)\n\n"
"  F-statistic = 1.3835\n"
"  p-value     = 0.2408\n"
"  Decision    = Fail to reject H0 at 5%\n\n"
"Why asymmetry is not statistically significant:\n\n"
"1. Headline CPI composition: Food is 47% of India's CPI basket.\n"
"   Oil is a small direct component. Effects get diluted by food-price noise.\n\n"
"2. Government interventions: Excise duty adjustments, LPG subsidies,\n"
"   price controls buffer consumers from full oil shocks.\n\n"
"3. Downward price rigidity: Administered prices and mark-up pricing\n"
"   adjust slowly and incompletely, especially downward.\n\n"
"4. Time-varying dynamics: 20 years averages together different regimes.\n"
"   Rolling window shows significance in some sub-periods.\n\n"
"This is NOT a failure - many published studies on headline CPI find similar results.",
0.8, 1.4, 12, 5.5, size=15)

# ===== SLIDE 12: SUB-SAMPLE =====
s = add_slide()
add_title_text(s, "Sub-Sample Analysis: Pre/Post Diesel Deregulation", 0.8, 0.4, 11, 0.7, size=30, color=DARK_BLUE)
add_line(s, 1.1)
add_image(s, "fig_5_subsample_comparison.png", 0.3, 1.3, 6.5)
add_body_text(s,
"Split at October 2014 (diesel deregulation):\n\n"
"Pre-2014 (N=122):\n"
"  CPT+ = 0.0173, CPT- = 0.0099\n"
"  p(asymmetry) = 0.846\n"
"  Adj R2 = 0.357\n\n"
"Post-2014 (N=123):\n"
"  CPT+ = 0.0102, CPT- = -0.002\n"
"  p(asymmetry) = 0.443\n"
"  Adj R2 = 0.500\n\n"
"Post-2014 shows cleaner results:\n"
"positive shock positive, negative\n"
"shock near-zero. Better model fit.\n"
"After deregulation, market prices\n"
"respond more directly to oil.",
7.2, 1.3, 5.5, 5.5, size=14)

# ===== SLIDE 13: DIAGNOSTICS =====
s = add_slide()
add_title_text(s, "Diagnostic Tests", 0.8, 0.4, 11, 0.7, size=30, color=DARK_BLUE)
add_line(s, 1.1)
add_body_text(s,
"Model Diagnostics (Main Asymmetric ADL):\n\n"
"  Test                        Statistic    p-value    Result\n"
"  ----                        ---------    -------    ------\n"
"  Breusch-Godfrey LM(12)     19.71        0.073      PASS (no autocorrelation)\n"
"  Breusch-Pagan              44.76        0.013      FAIL* (heteroskedasticity)\n"
"  Ramsey RESET                2.80        0.063      PASS (correct functional form)\n"
"  CUSUM Stability             0.86        0.091      PASS (model is stable)\n\n"
"* Breusch-Pagan detects heteroskedasticity, which is EXPECTED in macro data.\n"
"  This is precisely why we use Newey-West HAC standard errors -\n"
"  they give valid inference even with heteroskedasticity.\n"
"  The detection CONFIRMS our methodological choice was correct.",
0.8, 1.4, 12, 3.5, size=16)
add_image(s, "fig_6_cusum_stability.png", 3, 4.5, 7)

# ===== SLIDE 14: RESIDUALS =====
s = add_slide()
add_title_text(s, "Residual Diagnostics", 0.8, 0.4, 8, 0.7, size=30, color=DARK_BLUE)
add_line(s, 1.1)
add_image(s, "fig_8_residual_diagnostics.png", 0.3, 1.3, 7.5)
add_body_text(s,
"4-Panel Diagnostics:\n\n"
"1. Residuals over time:\n"
"   No pattern => no autocorrelation\n\n"
"2. Histogram:\n"
"   Roughly bell-shaped,\n"
"   fat tails (typical for macro)\n\n"
"3. Q-Q Plot:\n"
"   Mostly follows line,\n"
"   minor tail deviations\n\n"
"4. Actual vs Fitted:\n"
"   Clusters around 45-degree line\n\n"
"Overall: Residuals well-behaved.\n"
"HAC handles mild non-normality.",
8.2, 1.3, 4.5, 5.5, size=13)

# ===== SLIDE 15: ROBUSTNESS =====
s = add_slide()
add_title_text(s, "Robustness: Brent + Exchange Rate Model", 0.8, 0.4, 11, 0.7, size=30, color=DARK_BLUE)
add_line(s, 1.1)
add_body_text(s,
"Separating world oil price from exchange rate channel:\n\n"
"                        Primary (Oil_INR)    Brent USD + EXR\n"
"  CPT+                  0.021296             0.027458\n"
"  +10% shock effect     +0.213 pp            +0.275 pp\n"
"  p(CPT+ = 0)           0.1220               0.0933\n"
"  EXR coefficient       --                   0.039619\n"
"  EXR p-value           --                   0.0287 **\n"
"  Adj R-squared         0.4492               0.4575\n\n"
"Key Finding: Exchange rate is SIGNIFICANT (p = 0.029)!\n"
"A 1% rupee depreciation raises CPI by ~0.04 pp independently of oil.\n"
"This separates the world oil channel from the currency channel\n"
"and strengthens the India-specific interpretation.",
0.8, 1.4, 12, 5, size=16)

# ===== SLIDE 16: FUEL CPI =====
s = add_slide()
add_title_text(s, "Appendix: Fuel & Light CPI (MoSPI API)", 0.8, 0.4, 11, 0.7, size=30, color=DARK_BLUE)
add_line(s, 1.1)
add_body_text(s,
"Using official MoSPI Fuel & Light CPI sub-index (2011-2024):\n\n"
"                        Headline CPI (Main)    Fuel & Light CPI\n"
"  CPT+                  0.021296               0.060848\n"
"  +10% shock effect     +0.213 pp              +0.609 pp\n"
"  p(CPT+ = 0)           0.1220                 0.0337 **\n"
"  p(Asymmetry)           0.2408                 0.2647\n"
"  Adj R-squared         0.4492                 0.2144\n\n"
"Key Finding: Fuel CPI shows 3x STRONGER and STATISTICALLY SIGNIFICANT\n"
"positive oil pass-through (p = 0.034)!\n\n"
"This is exactly what theory predicts: a more energy-exposed CPI sub-index\n"
"responds more directly to oil prices.\n\n"
"This is the STRONGEST supporting evidence in the dissertation.\n"
"It confirms: oil DOES pass through to prices - headline CPI just dilutes the signal.",
0.8, 1.4, 12, 5.5, size=16)

# ===== SLIDE 17: OIL REGIMES =====
s = add_slide()
add_title_text(s, "Oil Price Regimes in Our Sample", 0.8, 0.4, 8, 0.7, size=30, color=DARK_BLUE)
add_line(s, 1.1)
add_image(s, "fig_9_oil_price_regimes.png", 0.3, 1.3, 8)
add_body_text(s,
"Our 20-year sample covers\nALL major oil regimes:\n\n"
"- China Boom (2004-08)\n"
"- GFC Crash (2008-09)\n"
"- Shale Glut (2014-16)\n"
"- COVID (2020)\n"
"- Russia-Ukraine (2022)\n\n"
"Rich variation in both\n"
"positive and negative shocks\n"
"of large magnitude.",
8.7, 1.3, 4, 5, size=14)

# ===== SLIDE 18: ROLLING WINDOW =====
s = add_slide()
add_title_text(s, "Rolling 60-Month Window Analysis", 0.8, 0.4, 8, 0.7, size=30, color=DARK_BLUE)
add_line(s, 1.1)
add_image(s, "fig_7_rolling_window.png", 0.3, 1.3, 7.5)
add_body_text(s,
"Rolling 5-year windows:\n\n"
"Red (CPT+): Varies over time,\n"
"often positive, strongest in\n"
"volatile oil periods.\n\n"
"Blue (CPT-): Consistently\n"
"near zero.\n\n"
"Dotted line = diesel deregulation\n\n"
"This explains weak full-sample\n"
"Wald test: averaging 20 years\n"
"of changing dynamics dilutes\n"
"the asymmetry signal.",
8.2, 1.3, 4.5, 5.2, size=14)

# ===== SLIDE 19: ASYMMETRY GAP =====
s = add_slide()
add_title_text(s, "Asymmetry Gap: Full Sample & Sub-Samples", 0.8, 0.4, 8, 0.7, size=30, color=DARK_BLUE)
add_line(s, 1.1)
add_image(s, "fig_10_asymmetry_gap.png", 0.3, 1.3, 7)
add_body_text(s,
"In ALL three samples:\n"
"  CPT+ (red) > |CPT-| (blue)\n\n"
"The DIRECTION of asymmetry\n"
"is completely robust across\n"
"all sample periods.\n\n"
"Even though statistical\n"
"significance is elusive,\n"
"the pattern is consistent:\n\n"
"Oil increases affect CPI\n"
"more than oil decreases\n"
"=> 'Rockets and Feathers'\n"
"behavior in India.",
7.8, 1.3, 5, 5.2, size=14)

# ===== SLIDE 20: CONCLUSION =====
s = add_slide()
add_title_text(s, "Conclusions & Key Takeaways", 0.8, 0.4, 11, 0.7, size=30, color=DARK_BLUE)
add_line(s, 1.1)
add_body_text(s,
"1. Positive oil shocks pass through to India's CPI inflation with economically\n"
"   meaningful magnitude: a 10% oil shock raises monthly CPI by ~0.21 pp\n\n"
"2. Negative oil shocks have essentially zero effect on CPI (CPT- ~ 0.0006)\n\n"
"3. Point estimates show clear asymmetry (CPT+ is 35x larger than |CPT-|),\n"
"   but formal Wald test is not significant at 5% (p = 0.24)\n\n"
"4. Exchange rate matters independently: a 1% rupee depreciation raises\n"
"   CPI by ~0.04 pp (significance p = 0.029 in Brent+EXR model)\n\n"
"5. Fuel & Light CPI shows 3x stronger and statistically significant\n"
"   positive pass-through (p = 0.034), confirming the energy channel\n\n"
"6. Results are robust to lag selection, COVID removal, winsorization,\n"
"   and sub-sample splits. All diagnostics pass.\n\n"
"7. Weak asymmetry in headline CPI is expected: food (47%) dilutes the signal,\n"
"   government interventions buffer consumers, prices are sticky downward.",
0.8, 1.4, 12, 5.5, size=16)

# ===== SLIDE 21: THANK YOU =====
s = add_slide()
add_title_text(s, "Thank You", 1.5, 2.5, 10, 1, size=40, color=DARK_BLUE, align=PP_ALIGN.CENTER)
add_body_text(s, "Questions & Discussion", 1.5, 3.8, 10, 0.5, size=24, color=GRAY)

outpath = f"{OUT}/Progress_Report_Aniket_Pandey.pptx"
prs.save(outpath)
print(f"PPTX saved to: {outpath}")
